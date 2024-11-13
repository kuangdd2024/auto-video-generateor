# pip install python-pptx pdf2image PyMuPDF
import os
import shutil
import subprocess
from pptx import Presentation
import fitz  # PyMuPDF

import comtypes.client
import pythoncom

from common_utils import get_savepath, chat


def ppt_to_images_windows(pptx_path, code_name):
    image_dir = get_savepath(code_name, 'image', mkdir_ok=True)
    # 初始化COM库
    pythoncom.CoInitialize()
    # 启动PowerPoint
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    # 打开PPT文件
    presentation = powerpoint.Presentations.Open(pptx_path)

    # 遍历所有幻灯片
    images = []
    for i, slide in enumerate(presentation.Slides):
        # 将幻灯片导出为图片
        image_path = f'{image_dir}/image_{i + 100}.jpg'
        slide.Export(image_path, "JPG")
        images.append(image_path)

    # 关闭PPT文件
    presentation.Close()
    # 关闭PowerPoint
    powerpoint.Quit()
    # 清理COM库
    pythoncom.CoUninitialize()
    return images


def ppt_to_texts_windows(pptx_path, code_name):
    text_dir = get_savepath(code_name, 'text', mkdir_ok=True)

    # 打开PPT文件
    prs = Presentation(pptx_path)

    # 遍历所有幻灯片
    texts = []
    for i, slide in enumerate(prs.slides):
        # 提取幻灯片的备注
        notes_slide = slide.notes_slide
        if notes_slide:
            notes_text = '\n'.join([paragraph.text for paragraph in notes_slide.notes_text_frame.paragraphs])
            notes_text.replace('\n', '。')
        else:
            notes_text = '嗯。'
        texts.append(notes_text)
        notes_path = f'{text_dir}/text_{i + 100}.txt'
        with open(notes_path, 'w', encoding='utf-8') as file:
            file.write(notes_text)
    return texts


def ppt_to_pdf(inpath, code_name):
    pdf_dir = get_savepath(code_name, 'pdf', mkdir_ok=True)
    if inpath.endswith('.pdf'):
        pdf_path = os.path.join(pdf_dir, 'document.pdf')
        shutil.copyfile(inpath, pdf_path)
        ppt_path = ''
        return ppt_path, pdf_path

    ppt_path = os.path.join(pdf_dir, 'document.pptx')
    shutil.copyfile(inpath, ppt_path)

    # Step 1: 使用libreoffice将PPT转换为PDF
    libreoffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
    # libreoffice_path = "libreoffice"

    pdf_path = os.path.join(pdf_dir, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")
    command = [
        libreoffice_path,
        "--headless",
        "--convert-to", "pdf",
        ppt_path,
        "--outdir", pdf_dir
    ]

    # 执行转换命令
    subprocess.run(command, check=True)
    print(f"PDF saved at: {pdf_path}")
    return ppt_path, pdf_path


def pdf_to_images(pdf_path, code_name):
    image_dir = get_savepath(code_name, 'image', mkdir_ok=True)

    # Step 2: 将PDF的每一页转换为图片
    image_format = "png"
    dpi = 150
    # 检查 PDF 文件是否存在
    if not os.path.isfile(pdf_path):
        raise FileNotFoundError(f"The file {pdf_path} does not exist.")

    # 打开 PDF 文件
    pdf_document = fitz.open(pdf_path)
    images = []
    # 遍历每一页
    for page_num in range(pdf_document.page_count):
        # 获取页面
        page = pdf_document[page_num]

        # 设置缩放比例
        zoom = dpi / 72  # 默认 PDF 分辨率是 72 DPI
        matrix = fitz.Matrix(zoom, zoom)

        # 渲染页面为图片
        pix = page.get_pixmap(matrix=matrix)

        # 保存图片
        image_path = f'{image_dir}/image_{page_num + 100}.{image_format}'
        pix.save(image_path)
        print(f"Saved: {image_path}")
        images.append(image_path)
    return images


def pdf_to_texts(pdf_path, code_name):
    text_dir = get_savepath(code_name, 'text', mkdir_ok=True)

    # Step 2: 将PDF的每一页转换为文字
    # 检查 PDF 文件是否存在
    if not os.path.isfile(pdf_path):
        raise FileNotFoundError(f"The file {pdf_path} does not exist.")

    # 打开 PDF 文件
    pdf_document = fitz.open(pdf_path)
    texts = []
    # 遍历每一页
    for page_num in range(pdf_document.page_count):
        # 获取页面
        page = pdf_document[page_num]

        text = page.get_text()  # 提取当前页的文本
        prompt_text = f'请用口述风格简明扼要讲解以下PPT内容，限制在60字以内，仅输出讲解内容。\nPPT内容：{text}。\n注意：请直接输出讲解内容。'
        note_text = chat(prompt_text)
        texts.append(note_text)
        note_path = f'{text_dir}/text_{page_num + 100}.txt'
        with open(note_path, 'w', encoding='utf-8') as file:
            file.write(note_text)

        print(f"Saved notes: {note_path}")
    return texts


def ppt_to_texts(ppt_path, code_name):
    text_dir = get_savepath(code_name, 'text', mkdir_ok=True)

    # Step 3: 使用python-pptx提取每页备注并保存为TXT文件
    ppt = Presentation(ppt_path)
    texts = []
    for i, slide in enumerate(ppt.slides):
        note_text = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""

        # 遍历幻灯片中的所有形状
        # texts_tmp = []
        # for shape in slide.shapes:
        #     if hasattr(shape, 'text'):
        #         # 提取形状中的文字
        #         texts_tmp.append(shape.text)
        # texts_tmp = '。'.join(texts_tmp)
        # note_text = note_text or texts_tmp or "嗯。"

        texts.append(note_text)
        note_path = f'{text_dir}/text_{i + 100}.txt'
        with open(note_path, 'w', encoding='utf-8') as file:
            file.write(note_text)

        print(f"Saved notes: {note_path}")
    return texts
